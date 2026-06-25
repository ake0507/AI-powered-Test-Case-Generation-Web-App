import io
from typing import Optional

from fastapi import UploadFile, HTTPException, status
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


async def extract_text_from_file(file: UploadFile) -> str:
    filename = file.filename or "uploaded_file"
    data = await file.read()
    lower_name = filename.lower()

    try:
        if lower_name.endswith('.pdf'):
            return _extract_text_from_pdf(data)
        if lower_name.endswith('.docx'):
            return _extract_text_from_docx(data)
        if lower_name.endswith('.pptx'):
            return _extract_text_from_pptx(data)
        if lower_name.endswith('.txt') or lower_name.endswith('.md'):
            return data.decode('utf-8', errors='replace')
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail={
                'error': 'FileParseError',
                'message': f'Unable to parse uploaded file: {exc}',
            },
        )

    raise HTTPException(
        status_code=status.HTTP_400_BAD_REQUEST,
        detail={
            'error': 'UnsupportedFileType',
            'message': 'Only .txt, .md, .docx, .pdf, and .pptx files are supported',
        },
    )


def _extract_text_from_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    text_chunks = []
    for page in reader.pages:
        page_text = page.extract_text() or ''
        text_chunks.append(page_text)
    return '\n\n'.join(text_chunks).strip()


def _extract_text_from_docx(data: bytes) -> str:
    document = Document(io.BytesIO(data))
    paragraphs = [para.text for para in document.paragraphs if para.text]
    return '\n\n'.join(paragraphs).strip()


def _extract_text_from_pptx(data: bytes) -> str:
    presentation = Presentation(io.BytesIO(data))
    slides = []
    for slide in presentation.slides:
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                slide_texts.append(shape.text)
        if slide_texts:
            slides.append('\n'.join(slide_texts))
    return '\n\n'.join(slides).strip()
